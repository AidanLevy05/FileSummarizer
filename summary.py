import pdfplumber
from pptx import Presentation
from collections import Counter
import spacy
import os

current_dir = os.path.dirname(os.path.abspath(__file__))
model_path = os.path.join(current_dir, "en_core_web_sm")
nlp = spacy.load(model_path)

def extract_text_from_pdf(pdf_path):
    text_per_slide = []
    try:
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                page_text = page.extract_text()
                if page_text:
                    text_per_slide.append(page_text)
    except Exception as e:
        print(f"Error reading PDF: {e}")
    return text_per_slide

def extract_text_from_pptx(pptx_path):
    text_per_slide = []
    try:
        presentation = Presentation(pptx_path)
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    slide_text.append(shape.text)
            text_per_slide.append("\n".join(slide_text))
    except Exception as e:
        print(f"Error reading PPTX: {e}")
    return text_per_slide

def extract_text_from_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            content = f.read()
        return [content] if content.strip() else []
    except Exception as e:
        print(f"Error reading TXT: {e}")
        return []

def summarize_text_spacy(text):
    doc = nlp(text)
    sentences = list(doc.sents)
    word_freq = Counter(token.text.lower() for token in doc if token.is_alpha and not token.is_stop)
    max_freq = max(word_freq.values()) if word_freq else 1
    for word in word_freq:
        word_freq[word] /= max_freq
    sentence_scores = {}
    for sent in sentences:
        sent_score = sum(word_freq.get(token.text.lower(), 0) for token in sent)
        if len(sent.text.split()) > 4:
            sent_score += 0.2 * len([ent for ent in sent.ents])
        sentence_scores[sent] = sent_score
    ranked_sentences = sorted(sentence_scores, key=sentence_scores.get, reverse=True)
    total_sentences = len(sentences)
    summary_count = max(3, total_sentences // 4)
    selected = ranked_sentences[:summary_count]
    selected = sorted(selected, key=lambda s: sentences.index(s))
    return " ".join(sent.text.strip() for sent in selected)

def process_file(file_path):
    slides_text = []

    if file_path.endswith(".pdf"):
        slides_text = extract_text_from_pdf(file_path)
    elif file_path.endswith(".pptx"):
        slides_text = extract_text_from_pptx(file_path)
    elif file_path.endswith(".txt"):
        slides_text = extract_text_from_txt(file_path)
    else:
        print("Unsupported file type.")
        return []

    if not slides_text:
        return []

    summaries = []
    for slide in slides_text:
        if slide.strip():
            summary = summarize_text_spacy(slide)
            summaries.append(summary)
    return summaries

def main(file_path):
    summaries = process_file(file_path)
    if not summaries:
        return []
    return summaries
